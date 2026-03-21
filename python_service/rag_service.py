import os
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from pydantic import BaseModel
import chromadb
from chromadb.config import Settings
import google.generativeai as genai
from pypdf import PdfReader
from pptx import Presentation
import uvicorn
import typing
from typing import List, Optional
import urllib.parse
import json

app = FastAPI()

# Paths
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
REPO_ROOT = os.path.dirname(BASE_DIR)
EDU_DIR = os.path.join(REPO_ROOT, "edu")
UPLOADS_DIR = os.path.join(EDU_DIR, "uploads", "syllabus")

# Configuration
GEMINI_API_KEY = "AIzaSyDoWzm5ec6PCr0K84P34yCpPnxa4lY3UY0"
genai.configure(api_key=GEMINI_API_KEY)

def get_model():
    # Use 'gemini-flash-latest' which maps to 1.5-flash and usually has better free tier quota
    # We saw 2.0-flash giving 429 'limit: 0' errors
    print("Initializing Gemini model (using gemini-flash-latest)...")
    try:
        name = 'models/gemini-flash-latest'
        m = genai.GenerativeModel(name)
        return m, name
    except Exception as e:
        print(f"Error initializing model: {str(e)}")
        # Fallback to 2.5-flash as another option from the list
        return genai.GenerativeModel('models/gemini-2.5-flash'), 'models/gemini-2.5-flash'

model, active_model_name = get_model()
print(f"ACTIVE MODEL: {active_model_name}")

# ChromaDB Setup - Disable telemetry to avoid PostHog connection errors
CHROMA_PATH = os.path.join(BASE_DIR, "chroma_db")
client = chromadb.PersistentClient(
    path=CHROMA_PATH,
    settings=Settings(anonymized_telemetry=False)
)
collection = client.get_or_create_collection(name="module_materials")

class QuestionRequest(BaseModel):
    moduleName: str
    classSubjectId: str
    difficulty: str = "MEDIUM"

class StudyPlanRequest(BaseModel):
    userId: int
    classSubjectId: int
    moduleName: Optional[str] = None
    days: int = 7
    stats: str

def extract_text(file_path: str):
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    if ext == ".pdf":
        reader = PdfReader(file_path)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    elif ext in [".ppt", ".pptx"]:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

@app.post("/ingest")
async def ingest_file(file_path: str, moduleName: str, classSubjectId: str):
    print(f"Ingest request: path={file_path}, module={moduleName}, id={classSubjectId}")
    
    # Decode URL-encoded paths (e.g., %20 -> space)
    file_path = urllib.parse.unquote(file_path)
    
    # Extract filename if it's a web path
    filename = os.path.basename(file_path)
    
    # Target path in repo structure
    final_path = os.path.join(UPLOADS_DIR, filename)
    
    print(f"Resolved path: {final_path}")
    
    if not os.path.exists(final_path):
        print(f"ERROR: File not found at {final_path}")
        raise HTTPException(status_code=404, detail=f"File not found: {final_path}")

    try:
        text = extract_text(final_path)
        if not text.strip():
            print("WARNING: Extracted text is empty")
            return {"status": "ignored", "reason": "empty text"}
        
        # Store in Chroma
        collection.add(
            documents=[text],
            metadatas=[{"moduleName": moduleName, "classSubjectId": classSubjectId}],
            ids=[f"{classSubjectId}_{moduleName}_{filename}"]
        )
        print("Success: Document indexed")
        return {"status": "success"}
    except Exception as e:
        print(f"ERROR during ingestion: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/generate-study-plan")
async def generate_study_plan(request: StudyPlanRequest):
    print(f"Generating study plan for user {request.userId}, subject {request.classSubjectId}")
    
    module_info = f"Module: {request.moduleName}" if request.moduleName and request.moduleName != "General" else "Overall Subject"
    
    prompt = f"""
    Student Performance Data for {module_info}:
    {request.stats}
    
    Please generate a very brief, encouraging, and personalized AI Study Plan for the next {request.days} days. 
    Highlight their weak modules if they scored below 50% and encourage re-reading the module materials. 
    Highlight strong areas and provide actionable next steps.
    
    IMPORTANT: Keep it to 2-3 sentences.
    """
    
    try:
        print(f"Using model: {active_model_name} for study plan")
        response = model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        print(f"ERROR during study plan generation: {str(e)}")
        # Try fallback if not already using it
        try:
            fallback_model = genai.GenerativeModel('models/gemini-2.0-flash')
            resp = fallback_model.generate_content(prompt)
            return resp.text.strip()
        except Exception as fe:
            print(f"Fallback also failed: {str(fe)}")
            raise HTTPException(status_code=500, detail=f"AI Generation Error: {str(e)}")

@app.post("/generate-questions")
async def generate_questions(request: QuestionRequest):
    # Query ChromaDB for relevant context
    results = collection.query(
        query_texts=[f"Generate questions for {request.moduleName}"],
        where={"$and": [{"moduleName": request.moduleName}, {"classSubjectId": request.classSubjectId}]},
        n_results=3
    )
    
    context = "\n".join(results['documents'][0]) if results['documents'] else ""
    
    if not context:
        raise HTTPException(status_code=404, detail="No context found for this module. Please ingest files first.")

    prompt = f"""
    Context: {context}
    
    Generate 5 multiple-choice questions (MCQs) and 2 short answer questions based ONLY on the context above.
    Difficulty: {request.difficulty}
    
    Output the response in STRICT JSON format:
    [
      {{
        "questionText": "...",
        "options": ["A", "B", "C", "D"],
        "correctOption": "...",
        "explanation": "...",
        "type": "MCQ",
        "difficulty": "{request.difficulty}"
      }},
      ...
    ]
    """
    
    try:
        # Prompt modified to handle non-JSON mode if needed
        full_prompt = prompt + "\n\nIMPORTANT: Output ONLY the JSON array. Do not include markdown formatting like ```json."
        
        print(f"Generating questions using model: {active_model_name}")
        response = model.generate_content(full_prompt)
        text = response.text.strip()
        
        # Clean up any potential markdown if the model ignores the instruction
        if "```" in text:
            if "```json" in text:
                text = text.split("```json")[1].split("```")[0].strip()
            else:
                text = text.split("```")[1].split("```")[0].strip()
        
        # Basic validation and parsing to ensure it's valid JSON for Java
        try:
            parsed_json = json.loads(text)
            print(f"Success: Generated {len(parsed_json)} questions for {request.moduleName}")
            return parsed_json
        except json.JSONDecodeError:
            print(f"ERROR: AI returned invalid JSON: {text[:200]}...")
            raise Exception("AI did not return valid JSON format.")
    except Exception as e:
        print(f"ERROR during generation with {active_model_name}: {str(e)}")
        try:
            fallback_name = 'models/gemini-flash-latest' if active_model_name != 'models/gemini-flash-latest' else 'models/gemini-pro-latest'
            print(f"Attempting fallback to {fallback_name}...")
            fallback_model = genai.GenerativeModel(fallback_name)
            resp = fallback_model.generate_content(prompt)
            return json.loads(resp.text)
        except Exception as fe:
            print(f"Fallback also failed: {str(fe)}")
            raise HTTPException(status_code=500, detail=f"AI Generation Error: {str(e)}")

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8001)
